# %% [markdown]
#  # สร้างไฟล์ PowerPoint ป้ายชื่อจากข้อมูลใน Excel
# %% [markdown]
#  ---
#  * author:  [Prasert Kanawattanachai](prasert.k@chula.ac.th)
#  * YouTube: https://www.youtube.com/prasertcbs
#  * github: https://github.com/prasertcbs/automate_pptx_nameplate
#  * [Chulalongkorn Business School](https://www.cbs.chula.ac.th/en/)
#  ---
# %% [markdown]
#  install required packages:
#  * pip: `pip install -U pandas Pillow qrcode python-pptx openpyxl`

# %%
import os
import sys
from io import BytesIO
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import qrcode  # https://pypi.python.org/pypi/qrcode
from PIL import Image  # https://pillow.readthedocs.io/en/stable/reference/Image.html


# %%
df=pd.read_excel('directory.xlsx', engine='openpyxl')
df


# %%
df[df['selected']]


# %%
def mecard(name, tel, email):
    """
    >>> mecard('Peter Parker', '088-123-4455', 'peter@marvel.com')
    'MECARD:N:Peter Parker;TEL:088-123-4455;EMAIL:peter@marvel.com;'
    """
    return f'MECARD:N:{name};TEL:{tel};EMAIL:{email};'

def gen_qrcode(text, logo_file=None, logo_pos='c'):
    qr = qrcode.QRCode(
        # About 30% or less errors can be corrected.
        error_correction=qrcode.constants.ERROR_CORRECT_M,
        border=1 # 10px
    )
    qr.add_data(text)
    qr.make(fit=True)
    qr_img = qr.make_image(fill_color="black", back_color="white")
    if logo_file is None:
        return qr_img
    else:
        logo_img = Image.open(logo_file).convert('RGBA')
        qr_logo_img = qr_img.copy().convert('RGBA')
        
        if logo_pos.lower() == 'c':  # center logo
            posx = (qr_img.size[0] - logo_img.size[0]) // 2
            posy = (qr_img.size[1] - logo_img.size[1]) // 2
        else:  # place logo on bottom right (margin 10px)
            margin=10
            posx = qr_img.size[0] - logo_img.size[0] - margin
            posy = qr_img.size[1] - logo_img.size[1] - margin

        qr_logo_img.paste(im=logo_img, box=(posx, posy), mask=logo_img)
        return qr_logo_img

def create_name_plate(pptx_output_name:str='directory.pptx', 
                      include_mecard_qr:bool=False, 
                      qr_logo:str='logo64x64.png'):
    """
    Summary: automate name plate creation in PowerPoint
    Usage:
    1. default: python nameplate.py
    2. specify pptx output file: python nameplate.py d1.pptx
    3. include mecard qr code: python nameplate.py d1.pptx y
    4. include mecard qr code with logo: python nameplate.py d1.pptx y logo64x64.png
    Args:
        pptx_output_name (str, optional): [description]. Defaults to 'directory.pptx'.
        include_mecard_qr (bool, optional): [description]. Defaults to False.
        qr_logo (str, optional): [description]. Defaults to 'logo64x64.png'.
    """
    df=pd.read_excel('directory.xlsx', engine='openpyxl')
    
    tmp_img_file_name='tmp_qrcode_img.png'
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts.get_by_name('Blank')
    for i, r in df[df['selected']].iterrows():
        slide = prs.slides.add_slide(blank_slide_layout)

        txBox = slide.shapes.add_textbox(left=Inches(1), top=Inches(3), 
                                         width=Inches(6), height=Inches(2))
        tf = txBox.text_frame
#         tf.word_wrap = False
#         tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tf.add_paragraph()
        p.font.bold=True
        p.font.size=Pt(80)
        p.font.color.rgb = RGBColor(128, 128, 128)
#         p.font.name='Browallia New'
        p.text = f"{r['fname']} {r['lname']}"

        if include_mecard_qr:
            qr_text=mecard(r['fname'] + ' ' + r['lname'], r['tel'], r['email'])
            img=gen_qrcode(qr_text, qr_logo, 'br')
            img.save(tmp_img_file_name)
            qr_left = Inches(7.5) 
            qr_top = Inches(3)
            qr_height=Inches(1.5)
            pic = slide.shapes.add_picture(tmp_img_file_name, 
                                           qr_left, qr_top, 
                                           height=qr_height)
    if include_mecard_qr:
        os.remove(tmp_img_file_name)
    prs.save(pptx_output_name)
    full_file_path=os.path.abspath(pptx_output_name)
    print(f'successfully saved file {full_file_path}')


# %%
# command line style
if __name__ == '__main__':
    print('''
    usage:
    1. default: python nameplate.py
    2. specify pptx output file: python nameplate.py d1.pptx
    3. include mecard qr code: python nameplate.py d1.pptx y
    4. include mecard qr code with logo: python nameplate.py d1.pptx y logo64x64.png
    ''')
    if len(sys.argv)==1:
        create_name_plate()
    elif len(sys.argv)==2:
        create_name_plate(sys.argv[1])
    elif len(sys.argv)==3:
        create_name_plate(sys.argv[1], True if sys.argv[2].lower()=='y' else False)
    elif len(sys.argv)==4:
        create_name_plate(sys.argv[1], True if sys.argv[2].lower()=='y' else False, sys.argv[3])
